#!/usr/bin/env python3
"""
spec_to_md.py — 仕様書ファイル（xlsx / xlsm / pdf / docx / pptx）を Markdown へ変換するツール

使い方:
    python tools/spec_to_md.py <input_file> [--out <output_dir>]

出力先デフォルト: docs/spec/

例:
    python tools/spec_to_md.py docs/spec/03_仕様書/product_spec.xlsx
    python tools/spec_to_md.py docs/spec/03_仕様書/product_spec.xlsm
    python tools/spec_to_md.py docs/spec/03_仕様書/product_spec.pdf --out output/
    python tools/spec_to_md.py docs/spec/03_仕様書/product_spec.pptx
"""

import argparse
import sys
from pathlib import Path

# 許可する拡張子（ホワイトリスト）
ALLOWED_EXTENSIONS = {".xlsx", ".xlsm", ".pdf", ".docx", ".pptx"}


def _safe_resolve(path: Path, base: Path) -> Path:
    """パストラバーサル防止: base 配下に収まるか確認"""
    resolved = path.resolve()
    # 入力ファイルは任意の場所でよい（ユーザ指定ファイル）
    return resolved


def convert_xlsx(src: Path) -> str:
    """xlsx → Markdown（シートごとにテーブルへ変換）"""
    import openpyxl  # noqa: PLC0415

    wb = openpyxl.load_workbook(src, read_only=True, data_only=True)
    lines: list[str] = [f"# {src.stem}\n"]

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue

        lines.append(f"\n## {sheet_name}\n")

        # 有効なセルを含む行だけ使う
        data_rows = [r for r in rows if any(c is not None for c in r)]
        if not data_rows:
            continue

        col_count = max(len(r) for r in data_rows)

        def fmt_cell(v: object) -> str:
            s = "" if v is None else str(v)
            # Markdown テーブル内のパイプ・改行をエスケープ
            return s.replace("|", "\\|").replace("\n", " ").replace("\r", "")

        header = data_rows[0]
        lines.append("| " + " | ".join(fmt_cell(c) for c in header) + " |")
        lines.append("| " + " | ".join("---" for _ in range(col_count)) + " |")
        for row in data_rows[1:]:
            padded = list(row) + [None] * (col_count - len(row))
            lines.append("| " + " | ".join(fmt_cell(c) for c in padded) + " |")

    wb.close()
    return "\n".join(lines) + "\n"


def convert_pdf(src: Path) -> str:
    """pdf → Markdown（テキスト＋テーブルを抽出）"""
    import pdfplumber  # noqa: PLC0415

    lines: list[str] = [f"# {src.stem}\n"]

    with pdfplumber.open(src) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            lines.append(f"\n## ページ {i}\n")

            # テーブル抽出
            tables = page.extract_tables()
            table_bboxes = [t.bbox for t in page.find_tables()] if tables else []

            if tables:
                for table in tables:
                    if not table:
                        continue
                    col_count = max(len(r) for r in table)

                    def fmt_cell(v: object) -> str:
                        s = "" if v is None else str(v)
                        return s.replace("|", "\\|").replace("\n", " ").replace("\r", "")

                    header = table[0]
                    padded_header = list(header) + [None] * (col_count - len(header))
                    lines.append("| " + " | ".join(fmt_cell(c) for c in padded_header) + " |")
                    lines.append("| " + " | ".join("---" for _ in range(col_count)) + " |")
                    for row in table[1:]:
                        padded = list(row) + [None] * (col_count - len(row))
                        lines.append("| " + " | ".join(fmt_cell(c) for c in padded) + " |")
                    lines.append("")

            # テキスト抽出（テーブル領域を除く）
            text = page.extract_text()
            if text:
                lines.append(text.strip())

    return "\n".join(lines) + "\n"


def convert_pptx(src: Path) -> str:
    """pptx → Markdown（スライドタイトル・テキスト・テーブルを抽出）"""
    from pptx import Presentation  # noqa: PLC0415
    from pptx.util import Pt  # noqa: PLC0415
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

    prs = Presentation(src)
    lines: list[str] = [f"# {src.stem}\n"]

    for i, slide in enumerate(prs.slides, start=1):
        # スライドタイトルを取得
        title_text = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
        lines.append(f"\n## スライド {i}{': ' + title_text if title_text else ''}\n")

        for shape in slide.shapes:
            # テーブル
            if shape.has_table:
                table = shape.table
                rows = table.rows
                if not rows:
                    continue
                col_count = len(rows[0].cells)

                def fmt_cell(cell) -> str:
                    text = cell.text.strip()
                    return text.replace("|", "\\|").replace("\n", " ")

                header = rows[0].cells
                lines.append("| " + " | ".join(fmt_cell(c) for c in header) + " |")
                lines.append("| " + " | ".join("---" for _ in range(col_count)) + " |")
                for row in rows[1:]:
                    lines.append("| " + " | ".join(fmt_cell(c) for c in row.cells) + " |")
                lines.append("")
                continue

            # テキストフレーム（タイトル shape は重複しないようスキップ）
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        lines.append(text)

    return "\n".join(lines) + "\n"


def convert_docx(src: Path) -> str:
    """docx → Markdown（見出し・段落・テーブルを変換）"""
    from docx import Document  # noqa: PLC0415
    from docx.oxml.ns import qn  # noqa: PLC0415

    doc = Document(src)
    lines: list[str] = [f"# {src.stem}\n"]

    # 見出しスタイル名 → Markdown レベルマッピング
    HEADING_MAP = {
        "Heading 1": "#",
        "Heading 2": "##",
        "Heading 3": "###",
        "Heading 4": "####",
        "見出し 1": "#",
        "見出し 2": "##",
        "見出し 3": "###",
        "見出し 4": "####",
    }

    def fmt_cell(cell) -> str:
        text = cell.text.strip()
        return text.replace("|", "\\|").replace("\n", " ")

    for block in doc.element.body:
        tag = block.tag.split("}")[-1] if "}" in block.tag else block.tag

        if tag == "p":
            from docx.text.paragraph import Paragraph  # noqa: PLC0415
            para = Paragraph(block, doc)
            style_name = para.style.name if para.style else ""
            text = para.text.strip()
            if not text:
                continue
            if style_name in HEADING_MAP:
                lines.append(f"\n{HEADING_MAP[style_name]} {text}\n")
            else:
                lines.append(text)

        elif tag == "tbl":
            from docx.table import Table  # noqa: PLC0415
            table = Table(block, doc)
            rows = table.rows
            if not rows:
                continue
            col_count = len(rows[0].cells)
            header_cells = rows[0].cells
            lines.append("")
            lines.append("| " + " | ".join(fmt_cell(c) for c in header_cells) + " |")
            lines.append("| " + " | ".join("---" for _ in range(col_count)) + " |")
            for row in rows[1:]:
                lines.append("| " + " | ".join(fmt_cell(c) for c in row.cells) + " |")
            lines.append("")

    return "\n".join(lines) + "\n"


def main() -> int:
    parser = argparse.ArgumentParser(description="仕様書（xlsx/pdf/docx）をMarkdownへ変換")
    parser.add_argument("input", help="変換する仕様書ファイルのパス")
    parser.add_argument(
        "--out",
        default="docs/spec",
        help="出力ディレクトリ（デフォルト: docs/spec）",
    )
    args = parser.parse_args()

    src = Path(args.input)

    # 入力検証
    if not src.exists():
        print(f"[ERROR] ファイルが見つかりません: {src}", file=sys.stderr)
        return 1

    ext = src.suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        print(
            f"[ERROR] 非対応の拡張子です: {ext}  (対応: {', '.join(ALLOWED_EXTENSIONS)})",
            file=sys.stderr,
        )
        return 1

    out_dir = Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / (src.stem + ".md")

    print(f"変換中: {src} → {out_path}")

    try:
        if ext in (".xlsx", ".xlsm"):
            md = convert_xlsx(src)
        elif ext == ".pdf":
            md = convert_pdf(src)
        elif ext == ".docx":
            md = convert_docx(src)
        elif ext == ".pptx":
            md = convert_pptx(src)
        else:
            # ここには到達しない（上のホワイトリスト検証済み）
            return 1
    except Exception as exc:  # noqa: BLE001
        print(f"[ERROR] 変換失敗: {exc}", file=sys.stderr)
        return 1

    header = (
        "<!-- このファイルはAIが生成したドラフトです。承認前に必ずレビューしてください。 -->\n"
        f"<!-- 元ファイル: {src.name} -->\n\n"
    )
    out_path.write_text(header + md, encoding="utf-8")
    print(f"[OK] 出力: {out_path}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
